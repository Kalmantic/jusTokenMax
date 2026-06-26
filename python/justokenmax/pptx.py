"""PPTX -> Markdown extraction.

A deck's text is already text — there is nothing to compress (unlike CSV/XLSX).
The win is *reliable* extraction: python-pptx recovers all slide text, including
text buried inside grouped shapes, renders tables, and lets us flag the
images/charts we drop so visual-only slides don't vanish silently. This mirrors
the DOCX handler (conversion, before == after tokens), not the XLSX one.

Thin extractor over python-pptx. The output conventions — titles as headings,
bold/italic runs, top-to-bottom/left-to-right reading order, merged-cell tables
— follow pptx2md (MIT, https://github.com/ssine/pptx2md), re-implemented rather
than vendored to avoid its heavy scipy/numpy dependency stack. Where pptx2md
extracts images to files, we instead flag them per slide: the goal is token
reduction, not asset export.
"""

from __future__ import annotations

from typing import List, Tuple

# Safety cap, same rationale as pdf.py — the Read hook runs this on untrusted
# files, so bound the work and the disk write.
MAX_OUTPUT_CHARS = 5_000_000

_TRUNCATED = "\n\n> _[justokenmax: output truncated — deck exceeds safety caps]_\n"


def _is_skippable_ph(shape) -> bool:
    """Slide-number / date / footer placeholders carry chrome (a page number, a
    date) not content; left in, they leak as noise bullets. Skip them."""
    from pptx.enum.shapes import PP_PLACEHOLDER
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type in (
                PP_PLACEHOLDER.SLIDE_NUMBER,
                PP_PLACEHOLDER.DATE,
                PP_PLACEHOLDER.FOOTER,
            )
    except Exception:
        return False
    return False


def _ph_type(shape):
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type
    except Exception:
        pass
    return None


def _is_picture(shape) -> bool:
    """True for a real picture — a PICTURE shape OR a *populated* picture
    placeholder (which reports shape_type PLACEHOLDER yet carries an image
    part). An empty picture placeholder has no image and is not a dropped
    visual, so it must not be flagged."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return True
    if getattr(shape, "is_placeholder", False):
        try:
            shape.image          # only a filled picture placeholder has this
            return True
        except Exception:
            return False
    return False


def _is_graphic_frame(shape) -> bool:
    """True for a graphicFrame element (SmartArt, OLE object, embedded media).
    Charts and tables are graphicFrames too but are handled by their own
    branches before the catch-all reaches this — so here it only matches the
    embedded objects we can't extract, which must still be flagged."""
    try:
        return shape._element.tag.endswith("}graphicFrame")
    except Exception:
        return False


def _runs_md(paragraph) -> str:
    """A paragraph's runs joined, with bold/italic mapped to Markdown. Boundary
    whitespace stays *outside* the emphasis markers — `**word** ` renders, but
    `**word **` (delimiter-adjacent space) does not in most Markdown."""
    out = []
    for run in paragraph.runs:
        t = run.text
        if not t:
            continue
        core = t.strip()
        if not core:
            out.append(t)                       # pure-whitespace run, keep as-is
            continue
        lead = t[:len(t) - len(t.lstrip())]
        trail = t[len(t.rstrip()):]
        if run.font.bold:
            core = f"**{core}**"
        if run.font.italic:
            core = f"_{core}_"
        out.append(lead + core + trail)
    return "".join(out)


def _para_lines(tf) -> List[str]:
    """Body paragraphs as bullets, indented by outline level."""
    lines = []
    for p in tf.paragraphs:
        txt = _runs_md(p).strip()
        if txt:
            lines.append("  " * p.level + f"- {txt}")
    return lines


def _cell(cell) -> str:
    # Spanned (non-origin) cells of a merge emit nothing — the origin holds the
    # text, so we never duplicate it across the merged span.
    if cell.is_spanned:
        return ""
    return cell.text.replace("|", "\\|").replace("\n", " ").strip()


def _table_md(table) -> str:
    rows = list(table.rows)
    if not rows:
        return ""
    out = []
    for ri, row in enumerate(rows):
        cells = [_cell(c) for c in row.cells]
        out.append("| " + " | ".join(cells) + " |")
        if ri == 0:
            out.append("| " + " | ".join("---" for _ in cells) + " |")
    return "\n".join(out)


def _pos(shape):
    """(top, left) in EMU for reading-order sorting. Fail-soft to (0, 0)."""
    try:
        return (shape.top or 0, shape.left or 0)
    except Exception:
        return (0, 0)


def _visual_note(images: int, charts: int, objects: int = 0) -> str:
    """Per-slide marker for dropped visuals (mirrors the xlsx/pdf markers).
    Empty when the slide has none. `objects` covers SmartArt/OLE/embedded
    graphic frames we can't extract — flagged so they never vanish silently."""
    parts = []
    if images:
        parts.append(f"{images} image(s)")
    if charts:
        parts.append(f"{charts} chart(s)")
    if objects:
        parts.append(f"{objects} embedded object(s)")
    if not parts:
        return ""
    return (f"> _[justokenmax: {' and '.join(parts)} on this slide not "
            "extracted — read the source if visual data matters]_")


def _walk(shapes, title_id, body: List[str], counters: dict) -> None:
    """Recurse shapes in reading order, appending body lines and counting the
    visuals we drop. Recurses into groups so nested text is never lost."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
    for sh in sorted(shapes, key=_pos):
        if sh.shape_id == title_id:
            continue                              # rendered as the slide heading
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            _walk(sh.shapes, title_id, body, counters)
            continue
        if _is_picture(sh):
            counters["images"] += 1
            continue
        if getattr(sh, "has_chart", False) and sh.has_chart:
            counters["charts"] += 1
            continue
        if getattr(sh, "has_table", False) and sh.has_table:
            md = _table_md(sh.table)
            if md:
                body.append(md)
                counters["tables"] += 1
            continue
        if _is_skippable_ph(sh):
            continue
        if sh.has_text_frame:
            if _ph_type(sh) == PP_PLACEHOLDER.SUBTITLE:
                sub = " ".join(sh.text_frame.text.split())
                if sub:
                    body.append(f"*{sub}*")
            else:
                lines = _para_lines(sh.text_frame)
                if lines:
                    body.append("\n".join(lines))
            continue
        # Catch-all: a graphicFrame that wasn't a chart or table is SmartArt, an
        # OLE object, or embedded media we can't extract — flag it so the visual
        # never disappears without a trace.
        if _is_graphic_frame(sh):
            counters["objects"] += 1


def _render_slide(slide, idx: int) -> Tuple[List[str], dict]:
    title_shape = slide.shapes.title
    # Collapse internal whitespace so a multi-paragraph title stays on the one
    # heading line instead of spilling a bare second line into the body.
    title = " ".join(title_shape.text.split()) if title_shape else ""
    title_id = title_shape.shape_id if title_shape else None

    body: List[str] = []
    counters = {"images": 0, "charts": 0, "tables": 0, "notes": 0, "objects": 0}
    _walk(slide.shapes, title_id, body, counters)

    # Drop page-number chrome typed into a plain text box: a lone bullet equal to
    # the slide's own number. Any other number (a KPI, a year) is real and kept —
    # we only strip the exact-position match, the one safe signal.
    page_noise = f"- {idx}"
    body = [b for b in body if b.strip() != page_noise]

    out = ["## Slide " + str(idx) + (f" — {title}" if title else "")]
    out.extend(body)

    if slide.has_notes_slide:
        nf = slide.notes_slide.notes_text_frame
        if nf is not None and nf.text.strip():
            out.append(f"> **Notes:** {nf.text.strip()}")
            counters["notes"] = 1

    marker = _visual_note(counters["images"], counters["charts"],
                          counters["objects"])
    if marker:
        out.append(marker)
    return out, counters


def pptx_to_markdown(path: str) -> Tuple[str, dict]:
    """Return (markdown, stats). stats = {slides, images, charts, tables, notes,
    objects}.

    Conversion, not compression: the text was already text, so optimize() bills
    this before == after. The value is reliable extraction + honest flagging.

    A failure loading the deck propagates (optimize() then fails open). But a
    single unrenderable slide must not lose the whole deck — it is flagged and
    the remaining slides still convert, so one exotic shape never silently turns
    a valid deck into a pass-through.
    """
    from pptx import Presentation

    prs = Presentation(path)
    parts: List[str] = []
    tot = {"slides": 0, "images": 0, "charts": 0, "tables": 0, "notes": 0,
           "objects": 0}
    for i, slide in enumerate(prs.slides, 1):
        try:
            lines, c = _render_slide(slide, i)
        except Exception:
            parts.append(f"## Slide {i}\n> _[justokenmax: slide {i} could not "
                         "be rendered — read the source]_")
            tot["slides"] += 1
            continue
        parts.append("\n".join(lines))
        tot["slides"] += 1
        for k in ("images", "charts", "tables", "notes", "objects"):
            tot[k] += c[k]

    md = "\n\n".join(parts).strip() + "\n"
    if len(md) > MAX_OUTPUT_CHARS:
        md = md[:MAX_OUTPUT_CHARS] + _TRUNCATED
    return md, tot
