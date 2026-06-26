"""PPTX -> Markdown extraction (optimize kind="pptx").

Builds real .pptx decks with python-pptx at runtime (no committed binaries,
matching the conftest convention) and asserts the slide-structured contract:
titles as headings, bullets with outline levels, bold/italic runs, Markdown
tables (incl. merged cells), shapes in top-to-bottom/left-to-right reading
order, text recovered from inside grouped shapes, and a per-slide marker for
dropped images/charts. Like DOCX, PPTX is conversion (before == after tokens),
not compression — the value is reliable extraction + honest visual flagging.
"""

import json

import pytest

# python-pptx is the optional `office` extra; skip (don't error) when absent.
pytest.importorskip("pptx")

from justokenmax.pptx import pptx_to_markdown
from justokenmax.optimize import optimize


def _img(path, color=(10, 80, 200)):
    from PIL import Image
    Image.new("RGB", (16, 16), color).save(str(path))
    return str(path)


# ---- text structure ----

def test_title_becomes_heading(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Roadmap"
    p = tmp_path / "t.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert "Roadmap" in md
    assert "## Slide 1" in md
    # the title shows up on the slide heading line, not as a body bullet
    heading = [ln for ln in md.splitlines() if ln.startswith("## Slide 1")][0]
    assert "Roadmap" in heading


def test_subtitle_is_classified(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = "Vision"
    slide.placeholders[1].text = "the one-liner"
    p = tmp_path / "s.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "Vision" in md
    assert "*the one-liner*" in md      # subtitle italicized, not a bullet


def test_body_bullets_preserve_levels(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Agenda"
    body = slide.placeholders[1].text_frame
    body.text = "Top"
    sub = body.add_paragraph()
    sub.text = "Nested"
    sub.level = 1
    p = tmp_path / "b.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "- Top" in md
    assert "  - Nested" in md            # one level of indent


def test_bold_and_italic_rendered(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    para = tf.paragraphs[0]
    r1 = para.add_run(); r1.text = "Important"; r1.font.bold = True
    r2 = para.add_run(); r2.text = " and ";
    r3 = para.add_run(); r3.text = "stressed"; r3.font.italic = True
    p = tmp_path / "fmt.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "**Important**" in md
    assert "_stressed_" in md


def test_multiline_title_stays_on_one_heading(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Line one\nLine two"
    p = tmp_path / "mlt.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    heading = [ln for ln in md.splitlines() if ln.startswith("## Slide 1")][0]
    assert "Line one" in heading and "Line two" in heading
    assert "\nLine two" not in md          # no stray second heading line


def test_bold_run_with_trailing_space_renders(tmp_path):
    # A bold run that ends in a space ("Bold ") must not produce "**Bold **",
    # which most Markdown renderers refuse to bold (delimiter-adjacent space).
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text_frame
    r1 = tf.paragraphs[0].add_run(); r1.text = "Bold "; r1.font.bold = True
    r2 = tf.paragraphs[0].add_run(); r2.text = "plain"
    p = tmp_path / "bs.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "**Bold**" in md
    assert "**Bold **" not in md


def test_nested_group_two_deep_is_extracted(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    g1 = slide.shapes.add_group_shape()
    g2 = g1.shapes.add_group_shape()
    g2.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)) \
        .text_frame.text = "deep-nested-text"
    p = tmp_path / "deep.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "deep-nested-text" in md


def test_one_unrenderable_slide_does_not_lose_the_deck(tmp_path, monkeypatch):
    # An exotic shape that crashes one slide must not abort the whole deck —
    # otherwise a single bad shape on slide 50 silently passes the entire binary
    # through (the maximal violation of "never silently drop").
    from justokenmax import pptx as m
    from pptx import Presentation

    prs = Presentation()
    for i in range(3):
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = f"S{i}"
    p = tmp_path / "three.pptx"
    prs.save(str(p))

    real = m._render_slide

    def flaky(slide, idx):
        if idx == 2:
            raise ValueError("exotic shape")
        return real(slide, idx)

    monkeypatch.setattr(m, "_render_slide", flaky)
    md, stats = m.pptx_to_markdown(str(p))
    assert "S0" in md and "S2" in md             # slides 1 and 3 survive
    assert "slide 2 could not be rendered" in md  # failure is flagged, not hidden
    assert stats["slides"] == 3


def test_graphic_frame_detector_recognises_non_text_frames(tmp_path):
    # SmartArt/OLE arrive as graphicFrame elements with neither chart nor table
    # APIs. The detector must recognise a graphicFrame so the catch-all flags it.
    # A chart's graphicFrame is a constructible specimen of the element type.
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from justokenmax.pptx import _is_graphic_frame

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData(); cd.categories = ["a"]; cd.add_series("s", (1,))
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(1), Inches(1), Inches(3), Inches(2), cd)
    tb = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(1))
    assert _is_graphic_frame(gf) is True
    assert _is_graphic_frame(tb) is False


def test_text_inside_group_is_extracted(tmp_path):
    # The recall path that silently regresses: a textbox nested in a group must
    # still be walked. Decks carry hundreds of groups, so this is load-bearing.
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    grp = slide.shapes.add_group_shape()
    grp.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)) \
        .text_frame.text = "buried-in-group"
    p = tmp_path / "g.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "buried-in-group" in md


def test_reading_order_top_to_bottom(tmp_path):
    # Shapes added out of visual order (lower one first) must render in
    # top-to-bottom order, not XML/insertion order.
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(5), Inches(3), Inches(1)) \
        .text_frame.text = "LOWER"
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)) \
        .text_frame.text = "UPPER"
    p = tmp_path / "order.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert md.index("UPPER") < md.index("LOWER")


# ---- tables ----

def test_table_rendered_as_markdown(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1),
                                 Inches(4), Inches(2)).table
    tbl.cell(0, 0).text = "Name"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "alpha"
    tbl.cell(1, 1).text = "42"
    p = tmp_path / "tbl.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert "| Name | Value |" in md
    assert "| alpha | 42 |" in md
    assert stats["tables"] == 1


def test_merged_cells_dont_duplicate(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tbl = slide.shapes.add_table(2, 2, Inches(1), Inches(1),
                                 Inches(4), Inches(2)).table
    a, b = tbl.cell(0, 0), tbl.cell(0, 1)
    a.merge(b)
    a.text = "Spanning"
    tbl.cell(1, 0).text = "x"
    tbl.cell(1, 1).text = "y"
    p = tmp_path / "merge.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))   # must not raise
    assert "Spanning" in md
    assert md.count("Spanning") == 1   # spanned cell not emitted twice


# ---- images / charts ----

def test_flags_images_per_slide(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(_img(tmp_path / "p.png"), Inches(1), Inches(1))
    p = tmp_path / "img.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert "1 image(s) on this slide not extracted" in md
    assert stats["images"] == 1


def test_image_count_is_exact(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(_img(tmp_path / "a.png", (255, 0, 0)),
                             Inches(1), Inches(1))
    slide.shapes.add_picture(_img(tmp_path / "b.png", (0, 0, 255)),
                             Inches(3), Inches(1))
    p = tmp_path / "imgs.pptx"
    prs.save(str(p))

    _, stats = pptx_to_markdown(str(p))
    assert stats["images"] == 2


def test_populated_picture_placeholder_is_flagged(tmp_path):
    # A picture dropped into a PICTURE placeholder reports shape_type PLACEHOLDER
    # (not PICTURE) and has no text frame — it must still be counted/flagged, or
    # template-based decks lose their images silently.
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[8])  # Picture with Caption
    ph = next(p for p in slide.placeholders
              if p.placeholder_format.type == PP_PLACEHOLDER.PICTURE)
    ph.insert_picture(_img(tmp_path / "pp.png"))
    p = tmp_path / "picph.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert stats["images"] == 1
    assert "1 image(s) on this slide not extracted" in md


def test_empty_picture_placeholder_is_not_flagged(tmp_path):
    # An *empty* picture placeholder is not a dropped image — flagging it would
    # be a false positive on every template slide with an unfilled frame.
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[8])  # empty picture placeholder
    p = tmp_path / "emptyph.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert stats["images"] == 0
    assert "image(s)" not in md


def test_flags_charts_per_slide(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cd = CategoryChartData()
    cd.categories = ["a", "b"]
    cd.add_series("S", (1, 2))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                           Inches(1), Inches(1), Inches(4), Inches(3), cd)
    p = tmp_path / "chart.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert stats["charts"] == 1
    assert stats["images"] == 0          # a chart is not an image
    assert "1 chart(s) on this slide not extracted" in md


# ---- notes ----

def test_speaker_notes_included(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Slide"
    slide.notes_slide.notes_text_frame.text = "remember to mention pricing"
    p = tmp_path / "notes.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))
    assert "remember to mention pricing" in md
    assert stats["notes"] == 1


# ---- noise placeholder filter ----

def test_slide_number_textbox_noise_is_dropped(tmp_path):
    # Authors often type page numbers into a plain text box (not a placeholder).
    # A lone text box equal to the slide's own number is chrome — drop it. A
    # different number is real data and must survive.
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    s1.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)) \
        .text_frame.text = "1"
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    s2.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)) \
        .text_frame.text = "2"
    s2.shapes.add_textbox(Inches(1), Inches(3), Inches(2), Inches(1)) \
        .text_frame.text = "42"
    p = tmp_path / "noise.pptx"
    prs.save(str(p))

    md, _ = pptx_to_markdown(str(p))
    assert "- 1" not in md      # slide 1's page-number "1" dropped
    assert "- 2" not in md      # slide 2's page-number "2" dropped
    assert "- 42" in md         # a real number is kept


def test_slide_number_placeholder_is_skippable():
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
    from justokenmax.pptx import _is_skippable_ph

    prs = Presentation()
    layout = prs.slide_layouts[5]
    sn = next(ph for ph in layout.placeholders
              if ph.placeholder_format.type == PP_PLACEHOLDER.SLIDE_NUMBER)
    body = next(ph for ph in prs.slide_layouts[1].placeholders
                if ph.placeholder_format.idx == 1)
    assert _is_skippable_ph(sn) is True
    assert _is_skippable_ph(body) is False


# ---- optimize() dispatch integration ----

def _deck(path, n=6):
    from pptx import Presentation

    prs = Presentation()
    for i in range(n):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Slide {i} title"
        slide.placeholders[1].text_frame.text = (
            f"Some body content for slide {i} with a few words to extract.")
    prs.save(str(path))


def test_optimize_detects_pptx(tmp_path):
    p = tmp_path / "deck.pptx"
    _deck(p)
    res = optimize(str(p))
    assert res.ok and res.kind == "pptx"
    assert res.output.endswith(".pptx.md")
    # Conversion, not compression: tokens are preserved (text was already text).
    assert res.tokens_after == res.tokens_before


def test_secret_in_pptx_is_masked(tmp_path):
    from pptx import Presentation

    secret = "AK" + "IA" + "S" * 16
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Creds"
    slide.placeholders[1].text_frame.text = f"aws key is {secret} keep safe"
    p = tmp_path / "leaky.pptx"
    prs.save(str(p))

    res = optimize(str(p))
    assert res.ok
    artifact = open(res.output, encoding="utf-8").read()
    assert secret not in artifact


def test_fail_open_when_library_missing(tmp_path, monkeypatch):
    from justokenmax import pptx as pptx_mod

    def _raise_import(_path):
        raise ImportError("No module named 'pptx'")

    monkeypatch.setattr(pptx_mod, "pptx_to_markdown", _raise_import)
    p = tmp_path / "deck.pptx"
    _deck(p)
    res = optimize(str(p))
    assert res.ok is False and res.kind == "skip"
    assert "python-pptx" in res.note


def test_real_corrupt_file_fails_open(tmp_path):
    # A genuinely broken .pptx (not a monkeypatched raise): python-pptx raises on
    # load, and optimize() must skip rather than crash the Read hook.
    p = tmp_path / "broken.pptx"
    p.write_bytes(b"PK\x03\x04 not actually a valid deck " * 200)
    res = optimize(str(p))
    assert res.ok is False and res.kind == "skip"


def test_cli_pptx_subcommand(tmp_path, capsys):
    from justokenmax.cli import main

    p = tmp_path / "deck.pptx"
    _deck(p)
    rc = main(["pptx", "--json", str(p)])
    out = json.loads(capsys.readouterr().out)
    assert rc == 0
    assert out["ok"] is True and out["kind"] == "pptx"


def test_output_size_cap_truncates(tmp_path, monkeypatch):
    from justokenmax import pptx as pptx_mod

    monkeypatch.setattr(pptx_mod, "MAX_OUTPUT_CHARS", 50)
    p = tmp_path / "deck.pptx"
    _deck(p, n=10)
    md, _ = pptx_to_markdown(str(p))
    assert "truncated" in md


def test_empty_slide_renders(tmp_path):
    from pptx import Presentation

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[6])  # blank, no shapes
    p = tmp_path / "empty.pptx"
    prs.save(str(p))

    md, stats = pptx_to_markdown(str(p))   # must not raise
    assert "## Slide 1" in md
    assert stats["slides"] == 1
